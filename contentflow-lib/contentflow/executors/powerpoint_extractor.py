"""PowerPoint extraction executor using python-pptx for text, slides, and images."""

import base64
import io
import logging
from pathlib import Path
from typing import Dict, Any, Optional, List

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    raise ImportError(
        "python-pptx is required for PowerPoint extraction. "
        "Install it with: pip install python-pptx"
    )

from . import ParallelExecutor
from ..models import Content

logger = logging.getLogger("contentflow.executors.powerpoint_extractor")


class PowerPointExtractorExecutor(ParallelExecutor):
    """
    Extract content from PowerPoint presentations using python-pptx.
    
    This executor analyzes PowerPoint presentations (.pptx) to extract text, slide-level chunks,
    tables, notes, and embedded images using the python-pptx library.
    
    Configuration (settings dict):
        - extract_text (bool): Extract full text content from presentation
          Default: True
        - extract_slides (bool): Create slide-level text chunks
          Default: True
        - extract_tables (bool): Extract tables from slides
          Default: True
        - extract_notes (bool): Extract speaker notes from slides
          Default: False
        - extract_properties (bool): Extract presentation properties
          Default: False
        - extract_images (bool): Extract embedded images from slides
          Default: False
        - content_field (str): Field containing PowerPoint bytes
          Default: None
        - temp_file_path_field (str): Field containing temp file path
          Default: "temp_file_path"
        - output_field (str): Field name for extracted data
          Default: "pptx_output"
        - image_output_mode (str): How to store extracted images
          Default: "base64"
          Options: "base64", "bytes"
        - slide_separator (str): Separator between slides in full text
          Default: "\n\n---\n\n"
        - include_slide_layouts (bool): Include slide layout information
          Default: False
          
        Also setting from ParallelExecutor and BaseExecutor apply.
    
    Example:
        ```python
        executor = PowerPointExtractorExecutor(
            id="pptx_extractor",
            settings={
                "extract_text": True,
                "extract_slides": True,
                "extract_tables": True,
                "extract_notes": True,
            }
        )
        ```
    
    Input:
        Document or List[Document] each with (ContentIdentifier) id containing:
        - data['content']: PowerPoint document bytes, OR
        - data['temp_file_path']: Path to PowerPoint file
        
    Output:
        Document or List[Document] with added fields:
        - data['pptx_output']['text']: Full extracted text
        - data['pptx_output']['slides']: List of slide chunks with text and metadata
        - data['pptx_output']['tables']: List of extracted tables (if enabled)
        - data['pptx_output']['notes']: List of speaker notes (if enabled)
        - data['pptx_output']['properties']: Presentation properties (if enabled)
        - data['pptx_output']['images']: List of extracted images (if enabled)
        - summary_data['slides_processed']: Number of slides processed
        - summary_data['tables_extracted']: Number of tables extracted
    """
    
    def __init__(
        self,
        id: str,
        settings: Optional[Dict[str, Any]] = None,
        **kwargs
    ):
        super().__init__(
            id=id,
            settings=settings,
            **kwargs
        )
        
        # Extract configuration
        self.extract_text = self.get_setting("extract_text", default=True)
        self.extract_slides = self.get_setting("extract_slides", default=True)
        self.extract_tables = self.get_setting("extract_tables", default=True)
        self.extract_notes = self.get_setting("extract_notes", default=False)
        self.extract_properties = self.get_setting("extract_properties", default=False)
        self.extract_images = self.get_setting("extract_images", default=False)
        self.content_field = self.get_setting("content_field", default=None)
        self.temp_file_field = self.get_setting("temp_file_path_field", default="temp_file_path")
        self.output_field = self.get_setting("output_field", default="pptx_output")
        self.image_output_mode = self.get_setting("image_output_mode", default="base64")
        self.slide_separator = self.get_setting("slide_separator", default="\n\n---\n\n")
        self.include_slide_layouts = self.get_setting("include_slide_layouts", default=False)
        
        # Validate image output mode
        if self.image_output_mode not in ["base64", "bytes"]:
            raise ValueError(f"Invalid image_output_mode: {self.image_output_mode}. Must be 'base64' or 'bytes'")
        
        if self.debug_mode:
            logger.debug(
                f"PowerPointExtractorExecutor with id {self.id} initialized: "
                f"text={self.extract_text}, slides={self.extract_slides}, "
                f"tables={self.extract_tables}, notes={self.extract_notes}"
            )
    
    async def process_content_item(
        self,
        content: Content
    ) -> Content:
        """Process a single PowerPoint presentation using python-pptx.
           Implements the abstract method from ParallelExecutor.
        """
        
        try:
            if not content or not content.data:
                raise ValueError("Content must have data")
            
            # Get PowerPoint content
            pptx_bytes = None
            pptx_path = None
            
            # Try to get from content field
            if self.content_field and self.content_field in content.data:
                pptx_bytes = content.data[self.content_field]
            
            # Try to get from temp file
            elif self.temp_file_field in content.data:
                pptx_path = content.data[self.temp_file_field]
            
            if not pptx_bytes and not pptx_path:
                raise ValueError(
                    f"PowerPoint missing required content. "
                    f"Needs either '{self.content_field}' or '{self.temp_file_field}'"
                )
            
            if self.debug_mode:
                source = f"file: {pptx_path}" if pptx_path else f"bytes: {len(pptx_bytes)} bytes"
                logger.debug(f"Processing PowerPoint {content.id} from {source}")
            
            prs = None
            
            try:
                # Open PowerPoint presentation
                if pptx_bytes:
                    prs = Presentation(io.BytesIO(pptx_bytes))
                else:
                    prs = Presentation(pptx_path)
            except Exception as e:
                logger.warning(
                    f"Invalid PowerPoint file for content {content.id}: {str(e)}"
                )
                content.summary_data['pptx_extraction_status'] = "invalid_file"
                return content
            
            extracted_data = {}
            
            # Extract slides
            all_text = []
            slides_data = []
            all_tables = []
            all_notes = []
            all_images = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                slide_tables = []
                
                # Extract text and tables from shapes
                for shape in slide.shapes:
                    # Extract text from shape
                    if hasattr(shape, "text") and shape.text:
                        slide_text.append(shape.text)
                    
                    # Extract tables
                    if self.extract_tables and shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        table_data = self._extract_table_from_shape(shape, slide_num + 1)
                        if table_data:
                            slide_tables.append(table_data)
                            all_tables.append(table_data)
                    
                    # Extract images
                    if self.extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_data = self._extract_image_from_shape(shape, slide_num + 1, len(all_images))
                        if image_data:
                            all_images.append(image_data)
                
                # Combine slide text
                slide_full_text = "\n".join(slide_text)
                
                if self.extract_text:
                    all_text.append(slide_full_text)
                
                if self.extract_slides:
                    slide_info = {
                        "slide_number": slide_num + 1,
                        "text": slide_full_text,
                        "char_count": len(slide_full_text),
                        "shape_count": len(slide.shapes),
                    }
                    
                    if self.include_slide_layouts:
                        slide_info["layout"] = slide.slide_layout.name if slide.slide_layout else None
                    
                    if self.extract_tables and slide_tables:
                        slide_info["tables"] = slide_tables
                    
                    slides_data.append(slide_info)
                
                # Extract notes
                if self.extract_notes:
                    notes_text = self._extract_notes_from_slide(slide, slide_num + 1)
                    if notes_text:
                        all_notes.append(notes_text)
            
            if self.extract_text:
                extracted_data['text'] = self.slide_separator.join(all_text)
                if self.debug_mode:
                    logger.debug(f"Extracted {len(extracted_data['text'])} characters of text")
            
            if self.extract_slides:
                extracted_data['slides'] = slides_data
                if self.debug_mode:
                    logger.debug(f"Created {len(slides_data)} slide chunks")
            
            if self.extract_tables:
                extracted_data['tables'] = all_tables
                if self.debug_mode:
                    logger.debug(f"Extracted {len(all_tables)} tables")
            
            if self.extract_notes:
                extracted_data['notes'] = all_notes
                if self.debug_mode:
                    logger.debug(f"Extracted {len(all_notes)} notes")
            
            if self.extract_images:
                extracted_data['images'] = all_images
                if self.debug_mode:
                    logger.debug(f"Extracted {len(all_images)} images")
            
            # Extract presentation properties
            if self.extract_properties:
                properties = self._extract_presentation_properties(prs)
                extracted_data['properties'] = properties
                if self.debug_mode:
                    logger.debug(f"Extracted presentation properties")
            
            # Store extracted data
            content.data[self.output_field] = extracted_data
            
            # Update summary
            content.summary_data['slides_processed'] = len(slides_data) if self.extract_slides else len(prs.slides)
            content.summary_data['tables_extracted'] = len(all_tables)
            content.summary_data['images_extracted'] = len(all_images)
            content.summary_data['extraction_status'] = "success"
                            
        except Exception as e:
            logger.error(
                f"PowerPointExtractorExecutor {self.id} failed processing presentation {content.id}",
                exc_info=True
            )
            
            # Raise the exception to be handled upstream if needed
            raise
        
        return content
    
    def _extract_table_from_shape(self, shape, slide_number: int) -> Optional[Dict[str, Any]]:
        """Extract table data from a shape.
        
        Args:
            shape: PowerPoint shape containing a table
            slide_number: Slide number where table is located
            
        Returns:
            Dictionary with table data or None if extraction fails
        """
        try:
            table = shape.table
            table_data = {
                "slide_number": slide_number,
                "rows": len(table.rows),
                "columns": len(table.columns),
                "data": []
            }
            
            # Extract cell data
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text)
                table_data["data"].append(row_data)
            
            return table_data
            
        except Exception as e:
            logger.warning(f"Failed to extract table from slide {slide_number}: {e}")
            return None
    
    def _extract_notes_from_slide(self, slide, slide_number: int) -> Optional[Dict[str, Any]]:
        """Extract speaker notes from a slide.
        
        Args:
            slide: PowerPoint slide
            slide_number: Slide number
            
        Returns:
            Dictionary with notes data or None if no notes
        """
        try:
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text and notes_text.strip():
                    return {
                        "slide_number": slide_number,
                        "text": notes_text,
                        "char_count": len(notes_text)
                    }
        except Exception as e:
            logger.warning(f"Failed to extract notes from slide {slide_number}: {e}")
        
        return None
    
    def _extract_image_from_shape(self, shape, slide_number: int, image_index: int) -> Optional[Dict[str, Any]]:
        """Extract image from a picture shape.
        
        Args:
            shape: PowerPoint picture shape
            slide_number: Slide number where image is located
            image_index: Index of the image
            
        Returns:
            Dictionary with image data or None if extraction fails
        """
        try:
            image = shape.image
            image_bytes = image.blob
            
            # Determine image format from content type
            content_type = image.content_type
            image_format = content_type.split('/')[-1] if '/' in content_type else 'unknown'
            
            # Prepare image data based on output mode
            if self.image_output_mode == "base64":
                image_data = base64.b64encode(image_bytes).decode('utf-8')
            else:
                image_data = image_bytes
            
            image_info = {
                "slide_number": slide_number,
                "image_index": image_index,
                "format": image_format,
                "size_bytes": len(image_bytes),
                "content_type": content_type,
                "data": image_data,
                "encoding": self.image_output_mode
            }
            
            return image_info
            
        except Exception as e:
            logger.warning(f"Failed to extract image from slide {slide_number}: {e}")
            return None
    
    def _extract_presentation_properties(self, prs: Presentation) -> Dict[str, Any]:
        """Extract presentation properties and metadata.
        
        Args:
            prs: Opened python-pptx Presentation
            
        Returns:
            Dictionary of presentation properties
        """
        properties = {}
        
        try:
            core_props = prs.core_properties
            
            properties['title'] = core_props.title or None
            properties['author'] = core_props.author or None
            properties['subject'] = core_props.subject or None
            properties['keywords'] = core_props.keywords or None
            properties['comments'] = core_props.comments or None
            properties['category'] = core_props.category or None
            properties['created'] = core_props.created.isoformat() if core_props.created else None
            properties['modified'] = core_props.modified.isoformat() if core_props.modified else None
            properties['last_modified_by'] = core_props.last_modified_by or None
            properties['revision'] = core_props.revision
            
            # Add presentation-specific info
            properties['slide_count'] = len(prs.slides)
            properties['slide_width'] = prs.slide_width
            properties['slide_height'] = prs.slide_height
            
        except Exception as e:
            logger.warning(f"Failed to extract some presentation properties: {e}")
        
        return properties
